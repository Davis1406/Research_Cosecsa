#!/usr/bin/env python3
"""
Convert PPTX slides to HTML for in-browser preview.
Usage: python3 pptx_to_html.py /path/to/file.pptx
Returns: JSON { slides: [...html strings], count: N }
"""
import sys, json, html, re, os

# Inject local packages for Apache's x86_64 Python (XAMPP runs under Rosetta on Apple Silicon)
# x86_64 wheels are installed in storage/python-packages-x86 via:
#   arch -x86_64 python3 -m pip install python-pptx --target storage/python-packages-x86
_here = os.path.dirname(os.path.abspath(__file__))
_app_root = os.path.dirname(_here)
_pkg_candidates = [
    os.path.join(_app_root, 'storage', 'python-packages-x86'),
    os.path.join(_app_root, 'storage', 'python-packages'),
]
for _d in _pkg_candidates:
    if os.path.isdir(_d) and _d not in sys.path:
        sys.path.append(_d)   # append (lowest priority) — don't shadow stdlib

try:
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print(json.dumps({"error": "python-pptx not installed"}))
    sys.exit(1)

SLIDE_W = 960
SLIDE_H = 540

def emu_to_px(emu, total_emu, total_px):
    if emu is None or total_emu == 0:
        return 0
    return round((emu / total_emu) * total_px)

def rgb_hex(color_obj):
    try:
        rgb = color_obj.rgb
        return "#{:02x}{:02x}{:02x}".format(rgb.red, rgb.green, rgb.blue)
    except:
        return None

def get_bg_color(slide, prs):
    """Try to extract slide background colour."""
    try:
        bg = slide.background
        fill = bg.fill
        fill.type  # trigger load
        if fill.type is not None:
            try:
                return "#{:02x}{:02x}{:02x}".format(
                    fill.fore_color.rgb.red,
                    fill.fore_color.rgb.green,
                    fill.fore_color.rgb.blue,
                )
            except:
                pass
    except:
        pass
    return "#ffffff"

def shape_to_html(shape, slide_w_emu, slide_h_emu):
    if not shape.has_text_frame:
        return ""

    left   = emu_to_px(shape.left,   slide_w_emu, SLIDE_W)
    top    = emu_to_px(shape.top,    slide_h_emu, SLIDE_H)
    width  = emu_to_px(shape.width,  slide_w_emu, SLIDE_W)
    height = emu_to_px(shape.height, slide_h_emu, SLIDE_H)

    paras_html = []
    for para in shape.text_frame.paragraphs:
        raw = para.text.replace('\x0b', '\n').strip()
        if not raw:
            paras_html.append('<div style="height:4px;"></div>')
            continue

        # Paragraph-level alignment
        align_map = {PP_ALIGN.CENTER: "center", PP_ALIGN.RIGHT: "right", PP_ALIGN.JUSTIFY: "justify"}
        text_align = align_map.get(para.alignment, "left")

        # Build run-level spans
        spans = []
        for run in para.runs:
            run_text = run.text.replace('\x0b', '\n').strip('\x0b')
            if not run_text:
                continue
            fs = 14
            bold   = False
            italic = False
            color  = None
            try:
                if run.font.size:
                    fs = max(8, min(54, int(run.font.size / 12700)))
                bold   = bool(run.font.bold)
                italic = bool(run.font.italic)
                if run.font.color and run.font.color.type:
                    color = rgb_hex(run.font.color)
            except:
                pass

            style = f"font-size:{fs}px;"
            if bold:   style += "font-weight:700;"
            if italic: style += "font-style:italic;"
            if color:  style += f"color:{color};"

            escaped = html.escape(run_text).replace('\n', '<br>')
            spans.append(f'<span style="{style}">{escaped}</span>')

        if not spans:
            # fallback: plain text
            escaped = html.escape(raw).replace('\n', '<br>')
            paras_html.append(f'<div style="text-align:{text_align}; margin:2px 0;">{escaped}</div>')
        else:
            paras_html.append(
                f'<div style="text-align:{text_align}; margin:2px 0; line-height:1.35;">{"".join(spans)}</div>'
            )

    if not paras_html:
        return ""

    return (
        f'<div style="position:absolute;left:{left}px;top:{top}px;'
        f'width:{width}px;min-height:{height}px;'
        f'overflow:hidden;padding:6px 8px;box-sizing:border-box;">'
        + "".join(paras_html)
        + '</div>'
    )

def slide_to_html(slide, prs, idx):
    try:
        slide_w_emu = prs.slide_width
        slide_h_emu = prs.slide_height
    except:
        slide_w_emu = 9144000
        slide_h_emu = 5143500

    bg = get_bg_color(slide, prs)
    shapes_html = []

    for shape in slide.shapes:
        try:
            s = shape_to_html(shape, slide_w_emu, slide_h_emu)
            if s:
                shapes_html.append(s)
        except:
            pass

    slide_html = (
        f'<div class="pptx-slide" data-slide="{idx+1}" '
        f'style="position:relative;width:{SLIDE_W}px;height:{SLIDE_H}px;'
        f'background:{bg};overflow:hidden;flex-shrink:0;">'
        + "".join(shapes_html)
        + '</div>'
    )
    return slide_html

def main():
    if len(sys.argv) < 2:
        print(json.dumps({"error": "No file path provided"}))
        sys.exit(1)

    path = sys.argv[1]
    try:
        prs = Presentation(path)
    except Exception as e:
        print(json.dumps({"error": f"Cannot open file: {str(e)}"}))
        sys.exit(1)

    slides = []
    for i, slide in enumerate(prs.slides):
        try:
            slides.append(slide_to_html(slide, prs, i))
        except Exception as e:
            slides.append(f'<div class="pptx-slide" data-slide="{i+1}" '
                          f'style="position:relative;width:{SLIDE_W}px;height:{SLIDE_H}px;'
                          f'background:#fff;display:flex;align-items:center;justify-content:center;">'
                          f'<span style="color:#aaa;font-size:14px;">Slide {i+1} — could not render</span></div>')

    print(json.dumps({"slides": slides, "count": len(slides)}))

if __name__ == "__main__":
    main()
